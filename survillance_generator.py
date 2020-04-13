from PyQt5.QtWidgets import QMainWindow, QApplication, QFileDialog, QLabel, QPushButton, QFrame
from PyQt5 import QtWidgets
from PyQt5 import QtCore
from PyQt5.QtGui import QPixmap, QImage
from PIL import Image, ImageQt
import Ui_main_window
import sys
from functools import partial
from pptx import Presentation
from pptx.util import Inches
import sip

class MyMainForm(QMainWindow, Ui_main_window.Ui_MainWindow):
    def __init__(self, parent=None):
        super(MyMainForm, self).__init__(parent)
        self.setupUi(self)
        self.pixMapHandler = DisplaySelectedImage()
        self.pixMapHandler.return_pixmapSignal.connect(self.ImageShowOnLabel)
        self.images = {}
        self.addImageWorkers = InsertImageToPPT()
        self.actiondaochu.triggered.connect(self.AddImageToPPT)
        self.ClearItems(self.pageLayout1)
        # self.AddEightGridLayoutToTab(self.pageLayout1, self.tab1)
        self.cb1.currentTextChanged.connect(partial(self.ChangeLayout, self.cb1))
    
    def ClearItems(self, layout):
        # layout.removeWidget(layout.children()[0])
        # sip.delete(layout.children()[0])
        items = (layout.itemAt(i) for i in range(layout.count()))
        for item in items:
            # item.widget().setParent(None)
            item.widget().deleteLater()
            # sip.delete(item)
    
    def ChangeLayout(self, cb):
        if cb == self.cb1:
            layout = self.pageLayout1
            page = self.tab1
        if layout.count() > 0:
            self.ClearItems(layout)
        if cb.currentText() == "八画面":
            self.AddEightGridLayoutToTab(layout, page)
        elif cb.currentText() == "十六画面":
            print("sixteen")
        else:
            print("twentyFive")
    
    def AddConnectToBtn(self, page):
        if self.images != None:
                pageNo = page.objectName()[-1]
                for key in list(self.images.keys()):
                    if key.split('-')[0] == pageNo:
                        del self.images[key]
        for btn in page.findChildren(QPushButton):
            no = btn.objectName().replace('pushButton', '')
            btn.clicked.connect(partial(self.InsertImage, no))
            self.images[no] = None
            

    def AddEightGridLayoutToTab(self, layout, page):
        pageNo = page.objectName()[-1]
        frame = QtWidgets.QFrame(self.horizontalLayoutWidget)
        frame.setFrameShape(QtWidgets.QFrame.StyledPanel)
        frame.setFrameShadow(QtWidgets.QFrame.Raised)
        frame.setObjectName("frame" + pageNo)
        frame_13 = QtWidgets.QFrame(frame)
        frame_13.setGeometry(QtCore.QRect(840, 340, 271, 161))
        frame_13.setFrameShape(QtWidgets.QFrame.Box)
        frame_13.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_13.setObjectName("frame" + pageNo + "_13")
        label1_4 = QtWidgets.QLabel(frame_13)
        label1_4.setGeometry(QtCore.QRect(10, 10, 251, 141))
        label1_4.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_4.setAlignment(QtCore.Qt.AlignCenter)
        label1_4.setObjectName("label" + pageNo + "_4")
        pushButton1_4 = QtWidgets.QPushButton(frame_13)
        pushButton1_4.setGeometry(QtCore.QRect(160, 130, 101, 23))
        pushButton1_4.setObjectName("pushButton" + pageNo + "_4")
        frame_15 = QtWidgets.QFrame(frame)
        frame_15.setGeometry(QtCore.QRect(840, 0, 271, 161))
        frame_15.setFrameShape(QtWidgets.QFrame.Box)
        frame_15.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_15.setObjectName("frame" + pageNo + "_15")
        label1_2 = QtWidgets.QLabel(frame_15)
        label1_2.setGeometry(QtCore.QRect(13, 4, 251, 151))
        label1_2.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_2.setAlignment(QtCore.Qt.AlignCenter)
        label1_2.setObjectName("label" + pageNo +"_2")
        pushButton1_2 = QtWidgets.QPushButton(frame_15)
        pushButton1_2.setGeometry(QtCore.QRect(160, 130, 101, 23))
        pushButton1_2.setObjectName("pushButton" + pageNo + "_2")
        frame_11 = QtWidgets.QFrame(frame)
        frame_11.setGeometry(QtCore.QRect(560, 510, 271, 161))
        frame_11.setFrameShape(QtWidgets.QFrame.Box)
        frame_11.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_11.setObjectName("frame" +pageNo + "_11")
        label1_7 = QtWidgets.QLabel(frame_11)
        label1_7.setGeometry(QtCore.QRect(10, 10, 251, 141))
        label1_7.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_7.setAlignment(QtCore.Qt.AlignCenter)
        label1_7.setObjectName("label"+ pageNo +"_7")
        pushButton1_7 = QtWidgets.QPushButton(frame_11)
        pushButton1_7.setGeometry(QtCore.QRect(160, 130, 101, 23))
        pushButton1_7.setObjectName("pushButton" + pageNo + "_7")
        frame_12 = QtWidgets.QFrame(frame)
        frame_12.setGeometry(QtCore.QRect(840, 510, 271, 161))
        frame_12.setFrameShape(QtWidgets.QFrame.Box)
        frame_12.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_12.setObjectName("frame" + pageNo + "_12")
        label1_8 = QtWidgets.QLabel(frame_12)
        label1_8.setGeometry(QtCore.QRect(10, 10, 251, 141))
        label1_8.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_8.setAlignment(QtCore.Qt.AlignCenter)
        label1_8.setObjectName("label" + pageNo + "_8")
        pushButton1_8 = QtWidgets.QPushButton(frame_12)
        pushButton1_8.setGeometry(QtCore.QRect(160, 130, 101, 23))
        pushButton1_8.setObjectName("pushButton" + pageNo + "_8")
        frame_14 = QtWidgets.QFrame(frame)
        frame_14.setGeometry(QtCore.QRect(840, 170, 271, 161))
        frame_14.setFrameShape(QtWidgets.QFrame.Box)
        frame_14.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_14.setObjectName("frame" + pageNo + "_14")
        label1_3 = QtWidgets.QLabel(frame_14)
        label1_3.setGeometry(QtCore.QRect(10, 10, 251, 141))
        label1_3.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_3.setAlignment(QtCore.Qt.AlignCenter)
        label1_3.setObjectName("label" + pageNo + "_3")
        pushButton1_3 = QtWidgets.QPushButton(frame_14)
        pushButton1_3.setGeometry(QtCore.QRect(160, 130, 101, 23))
        pushButton1_3.setObjectName("pushButton" + pageNo + "_3")
        frame_10 = QtWidgets.QFrame(frame)
        frame_10.setGeometry(QtCore.QRect(280, 510, 271, 161))
        frame_10.setFrameShape(QtWidgets.QFrame.Box)
        frame_10.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_10.setObjectName("frame" + pageNo + "_10")
        label1_6 = QtWidgets.QLabel(frame_10)
        label1_6.setGeometry(QtCore.QRect(10, 10, 251, 141))
        label1_6.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_6.setAlignment(QtCore.Qt.AlignCenter)
        label1_6.setObjectName("label" + pageNo + "_6")
        pushButton1_6 = QtWidgets.QPushButton(frame_10)
        pushButton1_6.setGeometry(QtCore.QRect(160, 130, 101, 23))
        pushButton1_6.setObjectName("pushButton" + pageNo + "_6")
        frame_9 = QtWidgets.QFrame(frame)
        frame_9.setGeometry(QtCore.QRect(0, 510, 271, 161))
        frame_9.setFrameShape(QtWidgets.QFrame.Box)
        frame_9.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_9.setObjectName("frame" + pageNo + "_9")
        label1_5 = QtWidgets.QLabel(frame_9)
        label1_5.setGeometry(QtCore.QRect(10, 10, 251, 141))
        label1_5.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_5.setAlignment(QtCore.Qt.AlignCenter)
        label1_5.setObjectName("label" + pageNo + "_5")
        pushButton1_5 = QtWidgets.QPushButton(frame_9)
        pushButton1_5.setGeometry(QtCore.QRect(160, 130, 101, 23))
        pushButton1_5.setObjectName("pushButton" + pageNo + "_5")
        frame_2 = QtWidgets.QFrame(frame)
        frame_2.setGeometry(QtCore.QRect(0, 0, 831, 501))
        frame_2.setFrameShape(QtWidgets.QFrame.Box)
        frame_2.setFrameShadow(QtWidgets.QFrame.Raised)
        frame_2.setObjectName("frame" + pageNo + "_2")
        label1_1 = QtWidgets.QLabel(frame_2)
        label1_1.setGeometry(QtCore.QRect(10, 10, 811, 481))
        label1_1.setFrameShape(QtWidgets.QFrame.NoFrame)
        label1_1.setAlignment(QtCore.Qt.AlignCenter)
        label1_1.setObjectName("label" + pageNo + "_1")
        pushButton1_1 = QtWidgets.QPushButton(frame_2)
        pushButton1_1.setGeometry(QtCore.QRect(720, 470, 101, 23))
        pushButton1_1.setObjectName("pushButton" + pageNo + "_1")
        _translate = QtCore.QCoreApplication.translate
        label1_4.setText(_translate("MainWindow", "图片4"))
        pushButton1_4.setText(_translate("MainWindow", "插入 / 更改图片"))
        label1_2.setText(_translate("MainWindow", "图片2"))
        pushButton1_2.setText(_translate("MainWindow", "插入 / 更改图片"))
        label1_7.setText(_translate("MainWindow", "图片7"))
        pushButton1_7.setText(_translate("MainWindow", "插入 / 更改图片"))
        label1_8.setText(_translate("MainWindow", "图片8"))
        pushButton1_8.setText(_translate("MainWindow", "插入 / 更改图片"))
        label1_3.setText(_translate("MainWindow", "图片3"))
        pushButton1_3.setText(_translate("MainWindow", "插入 / 更改图片"))
        label1_6.setText(_translate("MainWindow", "图片6"))
        pushButton1_6.setText(_translate("MainWindow", "插入 / 更改图片"))
        label1_5.setText(_translate("MainWindow", "图片5"))
        pushButton1_5.setText(_translate("MainWindow", "插入 / 更改图片"))
        label1_1.setText(_translate("MainWindow", "图片1"))
        pushButton1_1.setText(_translate("MainWindow", "插入 / 更改图片"))
        layout.addWidget(frame)
        self.AddConnectToBtn(page)
    
    def InsertImage(self, no):
        self.currentImagePath = QFileDialog.getOpenFileName(self, "选择图片文件","", "JPG 图片(*.jpg *.jpeg)")[0]
        self.currentDisplayLabel = "label" + no
        self.GetDisplayImage(no)

    def GetDisplayImage(self, no):
        self.pixMapHandler.imagePath = self.currentImagePath
        self.images[no] = self.currentImagePath
        self.pixMapHandler.start()
        print(self.images)
    
    def ImageShowOnLabel(self, pixMapImage):
        self.current_label = self.tab1.findChild((QLabel), self.currentDisplayLabel)
        self.current_label.setPixmap(pixMapImage)
        self.current_label.setScaledContents(True)
    
    def AddImageToPPT(self):
        self.addImageWorkers.form = self
        self.addImageWorkers.start()


class DisplaySelectedImage(QtCore.QThread):
    return_pixmapSignal = QtCore.pyqtSignal(QPixmap)

    def __init__(self):
        super(DisplaySelectedImage, self).__init__()
        self.imagePath = None
    
    def run(self):
        self.ConvertImageToQPixMap()

    def ConvertImageToQPixMap(self):
        if self.imagePath != None:
            pixMapImage = QPixmap(self.imagePath)
            self.return_pixmapSignal.emit(pixMapImage)

class InsertImageToPPT(QtCore.QThread):

    def __init__(self):
        super(InsertImageToPPT, self).__init__()
        self.form = None
    
    def run(self):
        # print(self.images)
        self.prs = Presentation(r'template.pptx')
        self.AddEightGridPPTSlide(self.prs, self.form.images)
        # for no in list(self.form.images.keys()):
        #     if no.split('_')[0] == '1':
        #         self.AddEightGridPPTSlide(self.prs, self.form.images[no])

    # def AddImageToPPT(self):
    #     pass

    def AddEightGridPPTSlide(self, prs, images):
        # img_paths_keys = [key for key in list(images.keys()) if key.split('_')[0] == '1']
        img_paths_keys = [key for key in list(images.keys()) if key.split('_')[0] == '1']
        img_paths_keys.sort()
        print(img_paths_keys)
        img = [images[key] if images[key] != None else 'dahua.jpg' for key in img_paths_keys]
        # print(img)
        left, top, width, height = Inches(0), Inches(0.32), Inches(7.52), Inches(4.41)
        prs.slides[0].shapes.add_picture(img[0], left, top, width, height)
        left, top, width, height = Inches(7.59), Inches(0.32), Inches(2.41), Inches(1.47)
        prs.slides[0].shapes.add_picture(img[1], left,top,width,height)
        left, top, width, height = Inches(7.59), Inches(1.87), Inches(2.41), Inches(1.375)
        prs.slides[0].shapes.add_picture(img[2], left, top, width, height)
        left, top, width, height = Inches(7.59), Inches(3.325), Inches(2.41), Inches(1.41)
        prs.slides[0].shapes.add_picture(img[3], left, top, width, height)
        left, top, width, height = Inches(7.59), Inches(4.81), Inches(2.41), Inches(1.44)
        prs.slides[0].shapes.add_picture(img[7], left, top, width, height)
        left, top, width, height = Inches(5.07), Inches(4.81), Inches(2.45), Inches(1.44)
        prs.slides[0].shapes.add_picture(img[6], left, top, width, height)
        left, top, width, height = Inches(2.54), Inches(4.81), Inches(2.45), Inches(1.44)
        prs.slides[0].shapes.add_picture(img[5], left, top, width, height)
        left, top, width, height = Inches(0), Inches(4.81), Inches(2.46), Inches(1.44)
        prs.slides[0].shapes.add_picture(img[4], left, top, width, height)
        prs.save('new_test.pptx')

    def AddSixteenGridPPTSlide(self):
        pass

    def AddTwentyFiveGridPPTSlide(self):
        pass


if __name__ == '__main__':
    app = QApplication(sys.argv)
    myWin = MyMainForm()
    myWin.show()
    sys.exit(app.exec_())
